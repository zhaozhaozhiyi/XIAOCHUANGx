#!/usr/bin/env python3
"""
Extract every shape on every slide of a .pptx into a JSON dump.

Usage:
    python extract_pptx.py <path/to/deck.pptx>            # prints to stdout
    python extract_pptx.py <path/to/deck.pptx> -o dump.json

The dump captures the *actual* state of the export — text content, position,
size, and per-run typography (font name, size, bold, italic, color). Use this
as the ground truth for the fidelity audit; do not trust the export script's
intent.

Coordinates are reported in inches (rounded to 3 decimals) so they're
human-readable when comparing against rails like CONTENT_MAX_Y = 6.70".
"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    sys.stderr.write(
        "python-pptx is required. Install with: pip install python-pptx\n"
    )
    sys.exit(2)


def emu_to_in(emu: int | None) -> float | None:
    if emu is None:
        return None
    return round(emu / 914400, 3)


def color_repr(color) -> str | None:
    """Best-effort color extraction. Returns hex string or None."""
    if color is None:
        return None
    try:
        # ColorFormat.type may be None when no explicit color is set.
        if color.type is None:
            return None
        rgb = color.rgb
        if rgb is None:
            return None
        return f"#{str(rgb).lower()}"
    except (AttributeError, ValueError, TypeError):
        return None


def extract_runs(text_frame) -> list[dict]:
    runs = []
    for para in text_frame.paragraphs:
        for run in para.runs:
            font = run.font
            runs.append({
                "text": run.text,
                "font": font.name,
                "size_pt": float(font.size.pt) if font.size is not None else None,
                "bold": bool(font.bold) if font.bold is not None else None,
                "italic": bool(font.italic) if font.italic is not None else None,
                # Color is independent of font name/size: a run can inherit
                # font from the theme yet set its own color. Color drift is
                # one of the things this audit needs to catch, so don't gate
                # the extraction on unrelated font attributes.
                "color": color_repr(font.color),
            })
    return runs


def extract_shape(shape) -> dict:
    data = {
        "name": shape.name,
        "shape_type": str(shape.shape_type) if shape.shape_type is not None else None,
        "left_in": emu_to_in(shape.left),
        "top_in": emu_to_in(shape.top),
        "width_in": emu_to_in(shape.width),
        "height_in": emu_to_in(shape.height),
    }
    if shape.left is not None and shape.height is not None and shape.top is not None:
        data["bottom_in"] = emu_to_in(shape.top + shape.height)
        data["right_in"] = emu_to_in(shape.left + shape.width)
    if shape.has_text_frame:
        tf = shape.text_frame
        data["text"] = tf.text
        data["runs"] = extract_runs(tf)
    return data


def extract_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    canvas = {
        "width_in": emu_to_in(prs.slide_width),
        "height_in": emu_to_in(prs.slide_height),
    }
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        shapes = [extract_shape(s) for s in slide.shapes]
        slides.append({"index": i, "shapes": shapes})
    return {
        "source": str(path),
        "canvas": canvas,
        "slide_count": len(slides),
        "slides": slides,
    }


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("path", type=Path, help=".pptx file to extract")
    ap.add_argument("-o", "--output", type=Path, help="write JSON to this path; default stdout")
    args = ap.parse_args()

    if not args.path.exists():
        ap.error(f"file not found: {args.path}")

    data = extract_pptx(args.path)
    payload = json.dumps(data, ensure_ascii=False, indent=2)
    if args.output:
        args.output.write_text(payload, encoding="utf-8")
        sys.stderr.write(f"wrote {args.output} ({len(payload)} bytes, {data['slide_count']} slides)\n")
    else:
        sys.stdout.write(payload)
        sys.stdout.write("\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
