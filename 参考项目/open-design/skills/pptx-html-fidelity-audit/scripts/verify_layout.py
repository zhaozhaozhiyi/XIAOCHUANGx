#!/usr/bin/env python3
"""
Verify a re-exported .pptx against footer-rail + canvas-bound invariants.

Usage:
    python verify_layout.py <path/to/deck.pptx>
    python verify_layout.py <path/to/deck.pptx> --content-max-y 6.70 --canvas-h 7.5

Exits 0 on no violations, 1 on any violation. Prints a single block of
violations sorted by slide index, one per line:

    slide 5  shape 'desc-row-B-1'  bottom 7.214" crosses footer rail 6.70"
    slide 11 shape 'note-paragraph' bottom 7.342" exceeds canvas 7.50"

Use this as the gate for "this re-export is shippable". Don't claim the audit
is fixed without running this script — the human eye misses 1–2 mm overflow
at zoom-out, the script doesn't.

Footer / chrome shapes are exempt from the content rail. Two heuristics
identify them, in this order:

1. **By name** — any shape whose name contains "footer", "foot", "chrome",
   "page", or "pagination" (case-insensitive). Use semantic names in your
   export script if you can.
2. **By position** — any shape whose `top` is at or below the footer-zone
   threshold (default `--footer-zone-top 6.80`). This catches python-pptx's
   auto-generated names like "TextBox 3" when the export script didn't name
   them. The threshold sits ~0.10" above FOOTER_TOP so chrome rows pinned
   exactly at FOOTER_TOP are still recognized.
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    sys.stderr.write(
        "python-pptx is required. Install with: pip install python-pptx\n"
    )
    sys.exit(2)


FOOTER_NAME_HINTS = ("footer", "foot", "chrome", "page", "pagination")
EPS_IN = 0.005   # ignore sub-pixel overflows (~0.13mm)


def is_footer_by_name(name: str) -> bool:
    n = (name or "").lower()
    return any(hint in n for hint in FOOTER_NAME_HINTS)


def emu_to_in(emu: int | None) -> float:
    return (emu or 0) / 914400


def verify(path: Path, content_max_y: float, canvas_w: float, canvas_h: float,
           footer_zone_top: float) -> list[str]:
    prs = Presentation(str(path))
    violations: list[str] = []

    actual_w = emu_to_in(prs.slide_width)
    actual_h = emu_to_in(prs.slide_height)
    if abs(actual_w - canvas_w) > EPS_IN or abs(actual_h - canvas_h) > EPS_IN:
        violations.append(
            f"canvas mismatch: file is {actual_w:.3f}\" x {actual_h:.3f}\", "
            f"expected {canvas_w}\" x {canvas_h}\""
        )

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.top is None or shape.height is None:
                continue
            top = emu_to_in(shape.top)
            left = emu_to_in(shape.left)
            bottom = top + emu_to_in(shape.height)
            right = left + emu_to_in(shape.width)
            name = shape.name or "<unnamed>"

            # Off-canvas (hard fail for any shape).
            if bottom > canvas_h + EPS_IN:
                violations.append(
                    f"slide {i:<2} shape '{name}' bottom {bottom:.3f}\" "
                    f"exceeds canvas {canvas_h}\""
                )
            if right > canvas_w + EPS_IN:
                violations.append(
                    f"slide {i:<2} shape '{name}' right {right:.3f}\" "
                    f"exceeds canvas width {canvas_w}\""
                )
            if top < -EPS_IN:
                violations.append(
                    f"slide {i:<2} shape '{name}' top {top:.3f}\" is negative"
                )
            if left < -EPS_IN:
                violations.append(
                    f"slide {i:<2} shape '{name}' left {left:.3f}\" is negative"
                )

            # Footer rail (only enforced on content shapes).
            # Shape is exempt if (a) named like a footer, or
            # (b) pinned at-or-below the footer zone threshold.
            if is_footer_by_name(name) or top >= footer_zone_top - EPS_IN:
                continue
            if bottom > content_max_y + EPS_IN:
                violations.append(
                    f"slide {i:<2} shape '{name}' bottom {bottom:.3f}\" "
                    f"crosses footer rail {content_max_y}\""
                )

    return violations


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    ap.add_argument("path", type=Path, help=".pptx file to verify")
    ap.add_argument("--content-max-y", type=float, default=6.70,
                    help="content rail in inches; nothing in content area may cross (default 6.70)")
    ap.add_argument("--canvas-w", type=float, default=13.333,
                    help="expected canvas width in inches (default 13.333 = 16:9)")
    ap.add_argument("--canvas-h", type=float, default=7.5,
                    help="expected canvas height in inches (default 7.5 = 16:9)")
    ap.add_argument("--footer-zone-top", type=float, default=6.80,
                    help="any shape with top >= this is treated as footer/chrome "
                         "(default 6.80; sits 0.10\" above the typical FOOTER_TOP=6.85\")")
    args = ap.parse_args()

    if not args.path.exists():
        ap.error(f"file not found: {args.path}")

    violations = verify(args.path, args.content_max_y, args.canvas_w, args.canvas_h,
                        args.footer_zone_top)
    if violations:
        sys.stderr.write("\n".join(violations) + "\n")
        sys.stderr.write(f"\n{len(violations)} violation(s) found in {args.path}\n")
        return 1
    sys.stderr.write(f"OK: 0 violations across all slides in {args.path}\n")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
